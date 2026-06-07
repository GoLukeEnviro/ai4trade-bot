from __future__ import annotations

import re
from html import escape
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

ROOT = Path(__file__).resolve().parents[1]
AUDIT_DIR = ROOT / "docs" / "audit"
REPORT_MD = AUDIT_DIR / "2026-05-29-ai4trade-bot-audit-report.md"
SLIDES_MD = AUDIT_DIR / "2026-05-29-ai4trade-bot-slide-outline.md"
PDF_OUT = AUDIT_DIR / "AI4Trade-Bot-Auditbericht-2026-05-29.pdf"
PPTX_OUT = AUDIT_DIR / "AI4Trade-Bot-Executive-Briefing-2026-05-29.pptx"


def _clean_inline_markdown(text: str) -> str:
    cleaned = text.strip()
    cleaned = cleaned.replace("**", "")
    cleaned = cleaned.replace("__", "")
    cleaned = cleaned.replace("`", "")
    cleaned = re.sub(r"\[(.*?)\]\((.*?)\)", r"\1", cleaned)
    return escape(cleaned)


def _build_styles():
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="AuditTitle",
            parent=styles["Title"],
            fontName="Helvetica-Bold",
            fontSize=22,
            leading=26,
            textColor=colors.HexColor("#0f172a"),
            spaceAfter=10,
        )
    )
    styles.add(
        ParagraphStyle(
            name="AuditH1",
            parent=styles["Heading1"],
            fontName="Helvetica-Bold",
            fontSize=16,
            leading=20,
            textColor=colors.HexColor("#0f172a"),
            spaceBefore=8,
            spaceAfter=6,
        )
    )
    styles.add(
        ParagraphStyle(
            name="AuditH2",
            parent=styles["Heading2"],
            fontName="Helvetica-Bold",
            fontSize=13,
            leading=17,
            textColor=colors.HexColor("#1d4ed8"),
            spaceBefore=6,
            spaceAfter=4,
        )
    )
    styles.add(
        ParagraphStyle(
            name="AuditBody",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=9.5,
            leading=13,
            textColor=colors.HexColor("#111827"),
            spaceAfter=4,
        )
    )
    styles.add(
        ParagraphStyle(
            name="AuditBullet",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=9.5,
            leading=13,
            textColor=colors.HexColor("#111827"),
            leftIndent=10,
            spaceAfter=3,
        )
    )
    return styles


def _add_page_chrome(canvas, doc):
    canvas.saveState()
    canvas.setFont("Helvetica", 8)
    canvas.setFillColor(colors.HexColor("#475569"))
    canvas.drawString(doc.leftMargin, 10 * mm, "AI4Trade Bot Audit | 2026-05-29")
    canvas.drawRightString(A4[0] - doc.rightMargin, 10 * mm, f"Seite {doc.page}")
    canvas.restoreState()


def build_pdf(report_path: Path, out_path: Path) -> int:
    styles = _build_styles()
    doc = SimpleDocTemplate(
        str(out_path),
        pagesize=A4,
        leftMargin=18 * mm,
        rightMargin=18 * mm,
        topMargin=18 * mm,
        bottomMargin=18 * mm,
        title="AI4Trade Bot Auditbericht 2026-05-29",
        author="GitHub Copilot",
    )

    story = []
    page_numbers: list[int] = []

    def on_page(canvas, doc):
        page_numbers.append(doc.page)
        _add_page_chrome(canvas, doc)

    for raw_line in report_path.read_text(encoding="utf-8").splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 3))
            continue
        if stripped.startswith("# "):
            story.append(Paragraph(_clean_inline_markdown(stripped[2:]), styles["AuditTitle"]))
            continue
        if stripped.startswith("## "):
            story.append(Paragraph(_clean_inline_markdown(stripped[3:]), styles["AuditH1"]))
            continue
        if stripped.startswith("### "):
            story.append(Paragraph(_clean_inline_markdown(stripped[4:]), styles["AuditH2"]))
            continue
        if stripped.startswith("- "):
            story.append(
                Paragraph(
                    _clean_inline_markdown(stripped[2:]),
                    styles["AuditBullet"],
                    bulletText="•",
                )
            )
            continue
        story.append(Paragraph(_clean_inline_markdown(stripped), styles["AuditBody"]))

    doc.build(story, onFirstPage=on_page, onLaterPages=on_page)
    return max(page_numbers) if page_numbers else 0


def parse_slides(slides_path: Path):
    deck_title = "AI4Trade Bot — Executive Briefing"
    slides: list[dict[str, list[str] | str]] = []
    current: dict[str, list[str] | str] | None = None

    for raw_line in slides_path.read_text(encoding="utf-8").splitlines():
        stripped = raw_line.strip()
        if not stripped:
            continue
        if stripped.startswith("# "):
            deck_title = stripped[2:].strip()
            continue
        if stripped.startswith("## Slide "):
            if current:
                slides.append(current)
            _, title = stripped.split(":", 1)
            current = {"title": title.strip(), "bullets": []}
            continue
        if stripped.startswith("- ") and current is not None:
            current["bullets"].append(stripped[2:].strip())

    if current:
        slides.append(current)

    return deck_title, slides


def build_pptx(slides_path: Path, out_path: Path) -> int:
    deck_title, slides = parse_slides(slides_path)
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = deck_title
    title_slide.placeholders[1].text = "Audit-Stichtag: 29. Mai 2026 | Quelle: Repository-Analyse"

    for slide_def in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(slide_def["title"])
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        bullets = list(slide_def["bullets"])[:5]
        for index, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(20 if index == 0 else 18)

    prs.save(str(out_path))
    return len(prs.slides)


def main() -> None:
    pdf_pages = build_pdf(REPORT_MD, PDF_OUT)
    slide_count = build_pptx(SLIDES_MD, PPTX_OUT)
    print(f"PDF erstellt: {PDF_OUT} ({pdf_pages} Seiten)")
    print(f"PPTX erstellt: {PPTX_OUT} ({slide_count} Folien)")


if __name__ == "__main__":
    main()
